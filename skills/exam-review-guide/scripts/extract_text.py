#!/usr/bin/env python3
"""Extract text from course materials (pptx, docx, pdf, md, txt) for exam review building.

Walks a directory, extracts text from every supported file, and writes one
Markdown file with per-file sections and per-slide/per-page labels so the
content stays traceable to its source.

Usage:
    python3 extract_text.py <directory> [-o output.md]

Dependencies (install on demand):
    pip install python-pptx python-docx pymupdf
    PDF fallback: the `pdftotext` command (poppler) is used if pymupdf is missing.

Legacy .ppt/.doc files are reported as unsupported; convert them first with
`soffice --headless --convert-to pptx/docx <file>` if LibreOffice is available.
"""

import argparse
import shutil
import subprocess
import sys
from pathlib import Path

SUPPORTED = {".pptx", ".docx", ".pdf", ".md", ".txt"}
LEGACY = {".ppt", ".doc"}


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    texts.append(" | ".join(cells))
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                texts.append(f"(备注: {note})")
        if texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


def extract_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx not installed. Run: pip install python-docx")
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            if para.style.name.startswith("Heading"):
                t = f"## {t}"
            parts.append(t)
    for table in doc.tables:
        rows = [" | ".join(c.text.strip() for c in row.cells) for row in table.rows]
        parts.append("[表格]\n" + "\n".join(rows))
    return "\n".join(parts)


def extract_pdf(path: Path) -> str:
    try:
        import fitz  # pymupdf
    except ImportError:
        return extract_pdf_fallback(path)
    parts = []
    with fitz.open(str(path)) as doc:
        for i, page in enumerate(doc, 1):
            t = page.get_text().strip()
            if t:
                parts.append(f"[Page {i}]\n{t}")
    return "\n\n".join(parts)


def extract_pdf_fallback(path: Path) -> str:
    if not shutil.which("pdftotext"):
        raise RuntimeError(
            "Neither pymupdf nor pdftotext is available. "
            "Run: pip install pymupdf  (or install poppler for pdftotext)"
        )
    result = subprocess.run(
        ["pdftotext", "-layout", str(path), "-"],
        capture_output=True, text=True, check=True,
    )
    pages = result.stdout.split("\f")
    parts = [f"[Page {i}]\n{p.strip()}" for i, p in enumerate(pages, 1) if p.strip()]
    return "\n\n".join(parts)


def extract_md(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace").strip()


EXTRACTORS = {
    ".pptx": extract_pptx,
    ".docx": extract_docx,
    ".pdf": extract_pdf,
    ".md": extract_md,
    ".txt": extract_md,
}


def main():
    ap = argparse.ArgumentParser(description="Extract text from pptx/docx/pdf/md/txt course files.")
    ap.add_argument("directory", help="directory to scan recursively")
    ap.add_argument("-o", "--output", default="extracted_text.md", help="output Markdown file")
    args = ap.parse_args()

    root = Path(args.directory).resolve()
    if not root.is_dir():
        sys.exit(f"Not a directory: {root}")

    # Skip the script's own output and generated review documents,
    # otherwise re-runs would ingest their own products.
    out_path = Path(args.output).resolve()
    skip_names = {"extracted_text.md"}

    def is_excluded(p: Path) -> bool:
        if p.resolve() == out_path:
            return True
        if p.name in skip_names:
            return True
        return p.name.startswith("复习资料") and p.suffix.lower() == ".md"

    files = sorted(
        p for p in root.rglob("*")
        if p.is_file()
        and p.suffix.lower() in SUPPORTED | LEGACY
        and not any(part.startswith(".") for part in p.relative_to(root).parts)
        and not is_excluded(p)
    )
    if not files:
        sys.exit("No supported files (.pptx/.docx/.pdf/.md/.txt/.ppt/.doc) found.")

    sections, failed = [], []
    for f in files:
        rel = f.relative_to(root)
        ext = f.suffix.lower()
        if ext in LEGACY:
            failed.append((rel, f"legacy {ext} format; convert with: soffice --headless --convert-to {ext[1:]}x"))
            continue
        try:
            text = EXTRACTORS[ext](f)
        except Exception as e:
            failed.append((rel, str(e)))
            continue
        if not text.strip():
            failed.append((rel, "no extractable text (may be image-only; needs OCR)"))
            continue
        sections.append(f"# FILE: {rel}\n\n{text}")
        print(f"ok      {rel}")

    for rel, reason in failed:
        print(f"FAILED  {rel}: {reason}", file=sys.stderr)

    out = Path(args.output)
    body = "\n\n---\n\n".join(sections)
    if failed:
        fail_list = "\n".join(f"- {rel}: {reason}" for rel, reason in failed)
        body += f"\n\n---\n\n# UNPARSED FILES\n\n{fail_list}"
    out.write_text(body, encoding="utf-8")
    print(f"\nWrote {out} ({len(sections)} parsed, {len(failed)} failed)")


if __name__ == "__main__":
    main()
